# -*- coding: utf-8 -*-
"""리서치 JSON payload -> 디자인된 PPTX 리포트 렌더러.

사용법: python build_pptx.py <payload.json> <output.pptx>

디자인 규칙 (고정):
- 포인트색 KB 옐로우 #FFBC00, 본문 다크그레이/화이트, 차분한 금융 리포트 톤
- 전체 폰트 '맑은 고딕' 단일 (latin/ea/cs 모두 지정 -> 글자 깨짐 방지)
- 슬라이드 순서: 표지 -> 종목 개요 -> 재무 요약 -> 가격/추세 -> 뉴스·심리 -> 리스크 -> 한 줄 종합
- 표는 데이터 7행 초과 시 자동 절단(생략 표기) -> 슬라이드 밖으로 잘리지 않음
"""
import json
import os
import sys
import tempfile

if hasattr(sys.stdout, "reconfigure"):  # Windows 콘솔(cp949) 한글 깨짐 방지
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FONT = "맑은 고딕"
ACCENT = RGBColor(0xFF, 0xBC, 0x00)   # KB 옐로우
DARK = RGBColor(0x26, 0x28, 0x2B)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MAX_DATA_ROWS = 7
DISCLAIMER = "본 리포트는 학습용 분석이며 투자 권유가 아닙니다."
COVER_NOTE = "무료 공개 데이터 기반 학습용 리포트입니다 · 투자 권유가 아닙니다"
BANNED = ["매수 추천", "매도 추천", "목표가", "목표주가", "적정주가",
          "사야 한다", "사야한다", "팔아야 한다", "팔아야한다"]


def set_font(run, size=12, bold=False, color=DARK):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) 또는 list of run-lists for 멀티런 문단."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.15
        p.space_after = Pt(6)
        runs = line if isinstance(line, list) else [line]
        for text, size, bold, color in runs:
            r = p.add_run()
            r.text = text
            set_font(r, size, bold, color)
    return tb


def bullet_lines(bullets, size=13):
    out = []
    for b in bullets:
        out.append([("▪ ", size, True, ACCENT), (str(b), size, False, DARK)])
    return out


def add_header(slide, title, page_no):
    add_rect(slide, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.34), ACCENT)
    add_text(slide, Inches(0.85), Inches(0.30), Inches(10.5), Inches(0.6),
             [(title, 21, True, DARK)])
    add_rect(slide, Inches(0.55), Inches(0.95), Inches(12.23), Pt(1.2), RGBColor(0xDD, 0xDD, 0xDD))
    # footer
    add_text(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.35),
             [("Stock-Team Research  ·  " + DISCLAIMER, 8, False, GRAY)])
    add_text(slide, Inches(12.2), Inches(7.05), Inches(0.6), Inches(0.35),
             [(str(page_no), 9, False, GRAY)], align=PP_ALIGN.RIGHT)


def add_source(slide, source, y=Inches(6.6)):
    if source:
        add_text(slide, Inches(0.55), y, Inches(12.2), Inches(0.4),
                 [("(" + source + ")", 9, False, GRAY)])


def make_chart_png(chart):
    """chart: {dates: [...], values: [...], title?, value_label?} -> 임시 PNG 경로."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams["font.family"] = "Malgun Gothic"
    plt.rcParams["axes.unicode_minus"] = False

    dates = chart["dates"]
    values = [float(v) for v in chart["values"]]
    step = max(1, len(dates) // 250)  # 다운샘플
    dates, values = dates[::step], values[::step]

    fig, ax = plt.subplots(figsize=(8.8, 3.5), dpi=150)
    x = range(len(values))
    ax.plot(x, values, color="#26282B", linewidth=1.6, zorder=3)
    ax.fill_between(x, values, min(values), color="#FFBC00", alpha=0.15, zorder=2)
    ax.scatter([len(values) - 1], [values[-1]], color="#FFBC00",
               edgecolor="#26282B", s=45, zorder=4)
    ax.annotate(f"{values[-1]:,.0f}", (len(values) - 1, values[-1]),
                textcoords="offset points", xytext=(-8, 10), ha="right",
                fontsize=9, fontweight="bold", color="#26282B")
    ticks = list(range(0, len(dates), max(1, len(dates) // 6)))
    ax.set_xticks(ticks)
    ax.set_xticklabels([str(dates[i]) for i in ticks], fontsize=8, color="#666666")
    ax.tick_params(axis="y", labelsize=8, colors="#666666")
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:,.0f}")
    if chart.get("title"):
        ax.set_title(chart["title"], fontsize=11, color="#26282B", loc="left", pad=8)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color("#CCCCCC")
    ax.grid(axis="y", color="#E8E8E8", linewidth=0.7, zorder=1)
    fig.tight_layout()
    path = os.path.join(tempfile.gettempdir(), "report_pptx_chart.png")
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def add_table(slide, table_spec, x, y, w):
    headers = table_spec["headers"]
    rows = [[str(c) for c in r] for r in table_spec.get("rows", [])]
    truncated = 0
    if len(rows) > MAX_DATA_ROWS:
        truncated = len(rows) - MAX_DATA_ROWS
        rows = rows[:MAX_DATA_ROWS]
    n_rows, n_cols = len(rows) + 1, len(headers)
    size = 12 if (n_cols <= 5 and n_rows <= 6) else (10 if n_cols <= 6 else 9)
    height = Inches(0.4 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, height)
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    # 첫 컬럼(항목명)을 약간 넓게
    first_w = int(w * (1.4 / (n_cols + 0.4)))
    rest_w = int((w - first_w) / max(1, n_cols - 1))
    for j in range(n_cols):
        tbl.columns[j].width = first_w if j == 0 else rest_w
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(h)
        set_font(r, size, True, DARK)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = val
            set_font(r, size, j == 0, DARK if j == 0 else GRAY)
    if truncated:
        add_text(slide, x, y + height + Inches(0.05), w, Inches(0.3),
                 [(f"※ 지면 관계상 {truncated}행 생략 — 상세는 원문 리서치 참조", 9, False, GRAY)])
    return y + height


def build(payload, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    name = payload["stock_name"]
    ticker = payload.get("ticker", "")
    date = payload["date"]

    # ── 1. 표지 ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT)
    add_text(s, Inches(1.1), Inches(2.1), Inches(11), Inches(0.5),
             [("STOCK-TEAM EQUITY RESEARCH", 13, True, GRAY)])
    title = f"{name} ({ticker})" if ticker else name
    add_text(s, Inches(1.1), Inches(2.6), Inches(11.5), Inches(1.4),
             [(title, 44, True, DARK)])
    add_rect(s, Inches(1.15), Inches(3.95), Inches(2.2), Pt(3), ACCENT)
    sub = payload.get("subtitle", "")
    lines = [(f"작성일  {date}", 16, False, GRAY)]
    if sub:
        lines.insert(0, (sub, 18, False, DARK))
    add_text(s, Inches(1.1), Inches(4.25), Inches(11), Inches(1.2), lines)
    add_text(s, Inches(1.1), Inches(6.8), Inches(11), Inches(0.4),
             [(COVER_NOTE, 11, True, GRAY)])

    # ── 2. 종목 개요 ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, "종목 개요", 2)
    ov = payload.get("overview", {})
    add_text(s, Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.8),
             bullet_lines(ov.get("bullets", []), 14))
    add_source(s, ov.get("source"))

    # ── 3. 재무 요약 ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, "재무 요약", 3)
    fin = payload.get("financials", {})
    bottom = Inches(1.35)
    if fin.get("table"):
        bottom = add_table(s, fin["table"], Inches(0.8), Inches(1.35), Inches(11.7))
    if fin.get("bullets"):
        add_text(s, Inches(0.8), bottom + Inches(0.45), Inches(11.8), Inches(2.0),
                 bullet_lines(fin["bullets"], 12))
    add_source(s, fin.get("source"))

    # ── 4. 가격/추세 ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, "가격 / 추세", 4)
    pr = payload.get("price", {})
    chart_y = Inches(1.3)
    img = None
    if pr.get("chart_image") and os.path.exists(pr["chart_image"]):
        img = pr["chart_image"]
    elif pr.get("chart"):
        img = make_chart_png(pr["chart"])
    if img:
        s.shapes.add_picture(img, Inches(0.8), chart_y, width=Inches(7.2))
        bx, bw = Inches(8.3), Inches(4.3)
    else:
        bx, bw = Inches(0.8), Inches(11.8)
    add_text(s, bx, chart_y + Inches(0.1), bw, Inches(4.6),
             bullet_lines(pr.get("bullets", []), 11.5 if img else 13))
    add_source(s, pr.get("source"))

    # ── 5. 뉴스·심리 ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, "뉴스 · 시장 심리", 5)
    nw = payload.get("news", {})
    add_text(s, Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.2),
             bullet_lines(nw.get("bullets", []), 12.5))
    if nw.get("sentiment"):
        add_rect(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.55), ACCENT)
        add_text(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.45),
                 [(f"시장 심리:  {nw['sentiment']}", 14, True, DARK)],
                 anchor=MSO_ANCHOR.MIDDLE)
    add_source(s, nw.get("source"))

    # ── 6. 리스크 ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, "핵심 리스크", 6)
    risks = payload.get("risks", [])[:3]
    y = Inches(1.35)
    box_h = Inches(1.62)
    for rk in risks:
        add_rect(s, Inches(0.8), y, Inches(11.7), box_h, LIGHT)
        add_rect(s, Inches(0.8), y, Inches(0.12), box_h, ACCENT)
        head = [[(rk.get("title", ""), 14, True, DARK),
                 (f"    영향: {rk.get('impact', '-')}", 11.5, True, GRAY)]]
        add_text(s, Inches(1.15), y + Inches(0.12), Inches(11.1), Inches(0.4), head)
        add_text(s, Inches(1.15), y + Inches(0.55), Inches(11.1), Inches(1.0),
                 [(rk.get("basis", ""), 11, False, GRAY)])
        y += box_h + Inches(0.18)
    add_source(s, payload.get("risks_source"), y=Inches(6.75))

    # ── 7. 한 줄 종합 ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, "종합", 7)
    sm = payload.get("summary", {})
    add_rect(s, Inches(1.4), Inches(2.6), Inches(0.16), Inches(1.5), ACCENT)
    add_text(s, Inches(1.85), Inches(2.4), Inches(10.0), Inches(2.0),
             [(sm.get("line", ""), 22, True, DARK)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.85), Inches(5.0), Inches(10.0), Inches(0.6),
             [(DISCLAIMER, 13, False, GRAY)])
    srcs = payload.get("data_sources", [])
    if srcs:
        lines = [("데이터 출처·기준일", 10, True, DARK)]
        lines += [("· " + str(x), 9, False, GRAY) for x in srcs[:6]]
        add_text(s, Inches(1.85), Inches(5.6), Inches(10.5), Inches(1.3), lines)
    add_source(s, sm.get("source"), y=Inches(6.95))

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    return len(prs.slides._sldIdLst)


def check_banned(payload):
    text = json.dumps(payload, ensure_ascii=False)
    hits = sorted({w for w in BANNED if w in text})
    if hits:
        print("[경고] 매수/매도·목표가 단정 표현 후보 감지:", ", ".join(hits))
        print("       인용 보도가 아니라 권유성 서술이라면 payload에서 제거 후 재실행하세요.")


def main():
    if len(sys.argv) != 3:
        print("사용법: python build_pptx.py <payload.json> <output.pptx>")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        payload = json.load(f)
    check_banned(payload)
    n = build(payload, sys.argv[2])
    print(f"OK: {sys.argv[2]} (슬라이드 {n}장)")


if __name__ == "__main__":
    main()
